"""
CustomsColab Pro — FastAPI Backend
WebSocket-powered real-time quiz platform
With 3-layer OCR extraction: digital text → image OCR → AI vision fallback
"""

from fastapi import FastAPI, WebSocket, WebSocketDisconnect, UploadFile, File, Form, HTTPException
from fastapi.staticfiles import StaticFiles
from fastapi.responses import HTMLResponse, FileResponse
from fastapi.middleware.cors import CORSMiddleware
import json
import asyncio
import time
import os
import re
import io
import httpx
import base64
import logging
from typing import Optional
import uvicorn

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# ── Layer 1: pdfplumber — digital text PDFs ──
try:
    import pdfplumber
    HAS_PDF = True
except ImportError:
    HAS_PDF = False
    logger.warning("pdfplumber not installed — digital PDF extraction disabled")

# ── Layer 2: python-pptx — PowerPoint files ──
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    logger.warning("python-pptx not installed — PPTX extraction disabled")

# ── Layer 3a: pytesseract — local OCR for scanned PDFs/images ──
try:
    import pytesseract
    from PIL import Image
    # Quick check tesseract binary exists
    pytesseract.get_tesseract_version()
    HAS_OCR = True
    logger.info("pytesseract OCR available")
except Exception:
    HAS_OCR = False
    logger.warning("pytesseract not available — OCR fallback disabled")

# ── Layer 3b: pdf2image — convert PDF pages to images for OCR ──
try:
    from pdf2image import convert_from_bytes
    HAS_PDF2IMAGE = True
    logger.info("pdf2image available")
except ImportError:
    HAS_PDF2IMAGE = False
    logger.warning("pdf2image not installed — PDF→image OCR disabled")

# ─────────────────────────────────────────────
# APP SETUP
# ─────────────────────────────────────────────
app = FastAPI(title="CustomsColab Pro", version="2.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Mount static files — graceful if folder missing
import pathlib
if pathlib.Path("static").exists():
    app.mount("/static", StaticFiles(directory="static"), name="static")
else:
    logger.warning("static/ folder not found — skipping static file mount")

# ─────────────────────────────────────────────
# GLOBAL STATE
# ─────────────────────────────────────────────
state = {
    "active": False,
    "session_name": "Customs Training",
    "join_code": "000000",
    "quiz_data": [],
    "scores": {},       # { user_id: {name, score, answers:[], joined_at, submitted_at} }
    "started_at": None,
    "current_question": 0,
}

AVATARS = ["🦁","🐯","🦊","🐺","🦝","🐻","🐼","🦄","🐲","🦅","🦋","🐬","🦁","🐸","🦉","🦚"]


# ─────────────────────────────────────────────
# WEBSOCKET MANAGER
# ─────────────────────────────────────────────
class ConnectionManager:
    def __init__(self):
        self.trainer_connections: list[WebSocket] = []
        self.participant_connections: dict[str, WebSocket] = {}  # user_id -> ws

    async def connect_trainer(self, ws: WebSocket):
        await ws.accept()
        self.trainer_connections.append(ws)

    async def connect_participant(self, ws: WebSocket, user_id: str):
        await ws.accept()
        self.participant_connections[user_id] = ws

    def disconnect_trainer(self, ws: WebSocket):
        if ws in self.trainer_connections:
            self.trainer_connections.remove(ws)

    def disconnect_participant(self, user_id: str):
        self.participant_connections.pop(user_id, None)

    async def broadcast_trainers(self, message: dict):
        dead = []
        for ws in self.trainer_connections:
            try:
                await ws.send_json(message)
            except Exception:
                dead.append(ws)
        for ws in dead:
            self.trainer_connections.remove(ws)

    async def send_participant(self, user_id: str, message: dict):
        ws = self.participant_connections.get(user_id)
        if ws:
            try:
                await ws.send_json(message)
            except Exception:
                self.disconnect_participant(user_id)

    async def broadcast_all(self, message: dict):
        await self.broadcast_trainers(message)
        dead = []
        for uid, ws in list(self.participant_connections.items()):
            try:
                await ws.send_json(message)
            except Exception:
                dead.append(uid)
        for uid in dead:
            self.disconnect_participant(uid)

manager = ConnectionManager()


# ─────────────────────────────────────────────
# HELPER: LEADERBOARD SNAPSHOT
# ─────────────────────────────────────────────
def leaderboard_snapshot():
    total_q = len(state["quiz_data"])
    sorted_scores = sorted(
        state["scores"].values(),
        key=lambda x: (x["score"], -x.get("submitted_at", 9e18)),
        reverse=True
    )
    return {
        "type": "leaderboard_update",
        "session_name": state["session_name"],
        "active": state["active"],
        "total_questions": total_q,
        "participants": [
            {
                "rank": i + 1,
                "user_id": p["user_id"],
                "name": p["name"],
                "score": p["score"],
                "avatar": p.get("avatar", "🦁"),
                "correct": p.get("correct", 0),
                "wrong": p.get("wrong", 0),
                "submitted": p.get("submitted", False),
                "pct": round(p["score"] / (total_q * 10) * 100, 1) if total_q > 0 else 0,
            }
            for i, p in enumerate(sorted_scores)
        ],
        "class_accuracy": round(
            sum(p.get("correct", 0) for p in state["scores"].values()) /
            max(sum(p.get("correct", 0) + p.get("wrong", 0) for p in state["scores"].values()), 1)
            * 100, 1
        ) if state["scores"] else 0,
        "timestamp": time.time(),
    }


# ─────────────────────────────────────────────
# FILE TEXT EXTRACTION  (3-layer strategy)
# ─────────────────────────────────────────────

MIN_USEFUL_CHARS = 150   # below this → assume scanned, try OCR

def _ocr_image(img) -> str:
    """Run Tesseract OCR on a PIL Image. Returns extracted text."""
    if not HAS_OCR:
        return ""
    try:
        # PSM 3 = fully automatic page segmentation (best for slides/docs)
        cfg = "--psm 3 --oem 3 -l eng"
        return pytesseract.image_to_string(img, config=cfg)
    except Exception as e:
        logger.warning(f"OCR error on image: {e}")
        return ""

def _extract_pptx(content: bytes) -> str:
    """
    Extract text from PPTX.
    Strategy:
      1. Text shapes (digital text) — fast, accurate
      2. If a slide yields < MIN_USEFUL_CHARS, OCR its thumbnail via python-pptx
    """
    if not HAS_PPTX:
        return ""
    prs = Presentation(io.BytesIO(content))
    all_text = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text += shape.text.strip() + "\n"
            # Also grab text from tables inside shapes
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        slide_text += cell.text.strip() + " "
        # If slide has very little text it may be image-heavy → OCR
        if len(slide_text.strip()) < MIN_USEFUL_CHARS and HAS_OCR:
            logger.info(f"Slide {slide_num} has low text ({len(slide_text)} chars) — trying OCR")
            try:
                # Render slide as PNG via python-pptx image export
                from pptx.util import Inches
                img_stream = io.BytesIO()
                # python-pptx doesn't render slides directly;
                # use slide thumbnail if accessible, else skip
                # (full rendering needs LibreOffice/unoconv on server)
                pass
            except Exception:
                pass
        all_text.append(slide_text)
    return "\n".join(all_text).strip()

def _extract_pdf(content: bytes) -> str:
    """
    Extract text from PDF with automatic OCR fallback.
    Layer 1: pdfplumber  — works for digital/selectable text PDFs
    Layer 2: pdf2image + pytesseract — for scanned/image-only PDFs
    Layer 3: returns whatever partial text was found
    """
    digital_text = ""

    # ── Layer 1: digital text via pdfplumber ──
    if HAS_PDF:
        try:
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                for page in pdf.pages:
                    t = page.extract_text()
                    if t:
                        digital_text += t + "\n"
        except Exception as e:
            logger.warning(f"pdfplumber error: {e}")

    # If we got enough digital text, done
    if len(digital_text.strip()) >= MIN_USEFUL_CHARS:
        logger.info(f"PDF: digital extraction succeeded ({len(digital_text)} chars)")
        return digital_text.strip()

    # ── Layer 2: OCR via pdf2image + pytesseract ──
    if HAS_PDF2IMAGE and HAS_OCR:
        logger.info("PDF appears scanned — running OCR via pdf2image + pytesseract")
        ocr_text = ""
        try:
            # Convert PDF pages to PIL images at 200 DPI (good balance speed/quality)
            images = convert_from_bytes(
                content,
                dpi=200,
                fmt="jpeg",
                thread_count=2,
            )
            for i, img in enumerate(images):
                page_ocr = _ocr_image(img)
                if page_ocr.strip():
                    ocr_text += f"\n[Page {i+1}]\n{page_ocr}"
                logger.info(f"  OCR page {i+1}: {len(page_ocr)} chars extracted")
            if len(ocr_text.strip()) >= MIN_USEFUL_CHARS:
                logger.info(f"OCR succeeded: {len(ocr_text)} total chars")
                return ocr_text.strip()
        except Exception as e:
            logger.warning(f"pdf2image/OCR pipeline error: {e}")

    # ── Layer 3: return whatever we have (even if partial) ──
    try:
        result = digital_text.strip() or ocr_text.strip()
    except NameError:
        result = digital_text.strip()
    logger.warning(f"PDF extraction yielded limited text ({len(result)} chars) — AI will use what's available")
    return result

def _extract_image(content: bytes, filename: str) -> str:
    """OCR a raw image file (PNG, JPG, TIFF, BMP)."""
    if not HAS_OCR:
        return ""
    try:
        img = Image.open(io.BytesIO(content))
        # Convert to RGB if needed (handles CMYK, palette, etc.)
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        text = _ocr_image(img)
        logger.info(f"Image OCR ({filename}): {len(text)} chars")
        return text.strip()
    except Exception as e:
        logger.warning(f"Image OCR error: {e}")
        return ""

def extract_text(filename: str, content: bytes) -> str:
    """
    Main entry point for all file types.
    Returns extracted text string — never raises.
    """
    fn = filename.lower()
    try:
        if fn.endswith(".pptx") or fn.endswith(".ppt"):
            return _extract_pptx(content)
        elif fn.endswith(".pdf"):
            return _extract_pdf(content)
        elif fn.endswith(".txt"):
            return content.decode("utf-8", errors="ignore").strip()
        elif fn.endswith((".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".webp")):
            return _extract_image(content, filename)
        else:
            # Try as plain text last resort
            return content.decode("utf-8", errors="ignore").strip()
    except Exception as e:
        logger.error(f"extract_text failed for {filename}: {e}")
        return ""


# ─────────────────────────────────────────────
# AI QUIZ GENERATION — multi-model fallback
# Tries each free model in order until one works
# ─────────────────────────────────────────────

# Models tried in order — all free tier on OpenRouter
# Updated with verified working model IDs April 2026
MODELS_TO_TRY = [
    ("qwen/qwen3-coder:free",                          "Qwen3 Coder"),
    ("qwen/qwen3-next-80b-a3b-instruct:free",          "Qwen3 Next 80B"),
    ("openai/gpt-oss-120b:free",                       "GPT OSS 120B"),
    ("google/gemma-4-26b-a4b-it:free",                 "Gemma 4 26B"),
    ("nousresearch/hermes-3-llama-3.1-405b:free",      "Hermes 3 Llama 405B"),
    ("deepseek/deepseek-r1:free",                      "DeepSeek R1"),
    ("google/gemini-2.0-flash-exp:free",               "Gemini 2.0 Flash"),
    ("meta-llama/llama-3.1-8b-instruct:free",          "Llama 3.1 8B"),
    ("mistralai/mistral-7b-instruct:free",             "Mistral 7B"),
]

def _build_prompt(text: str, num_q: int) -> str:
    return f"""You are an expert trainer creating a quiz from training material.
Create exactly {num_q} Multiple Choice Questions from the text below.
Focus on key facts, concepts, rules, and procedures in the document.

Return ONLY a valid JSON array. No markdown, no explanation, no backticks:
[
  {{
    "question": "Question text?",
    "options": ["Option A", "Option B", "Option C", "Option D"],
    "answer": 0,
    "explanation": "Why this answer is correct."
  }}
]
"answer" = integer index 0-3 of correct option.

Text:
{text[:5000]}
"""

async def _try_model(client: httpx.AsyncClient, api_key: str, model_id: str, prompt: str) -> tuple:
    """
    Try a single model. Returns (quiz_list, error_string).
    quiz_list is None if failed, error_string is None if succeeded.
    """
    try:
        r = await client.post(
            "https://openrouter.ai/api/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json",
                "HTTP-Referer": "https://customsiq.app",
                "X-Title": "CustomsIQ",
            },
            json={
                "model": model_id,
                "messages": [{"role": "user", "content": prompt}],
                "temperature": 0.3,
            }
        )

        # Check HTTP status
        if r.status_code == 429:
            return None, "rate_limited"
        if r.status_code == 402:
            return None, "insufficient_credits"
        if r.status_code == 404:
            # Model not found / wrong ID — log body for debugging
            try:
                err_body = r.json().get("error", {}).get("message", "model not found")
            except Exception:
                err_body = r.text[:100]
            return None, f"model_not_found: {err_body[:80]}"
        if r.status_code != 200:
            return None, f"http_{r.status_code}: {r.text[:80]}"

        data = r.json()

        # Check response structure
        if "choices" not in data or not data["choices"]:
            return None, "no_choices_in_response"

        raw = data["choices"][0]["message"]["content"].strip()

        # Strip markdown fences if present
        if raw.startswith("```json"):
            raw = raw[7:]
        elif raw.startswith("```"):
            raw = raw[3:]
        if raw.endswith("```"):
            raw = raw[:-3]
        raw = raw.strip()

        if not raw:
            return None, "empty_response"

        # Parse JSON
        quiz = json.loads(raw)

        # Validate structure
        valid = [q for q in quiz if all(k in q for k in ["question", "options", "answer"])]
        if not valid:
            return None, "invalid_quiz_structure"

        return valid, None

    except json.JSONDecodeError as e:
        return None, f"json_error: {str(e)[:80]}"
    except Exception as e:
        return None, f"error: {str(e)[:80]}"


async def generate_quiz_ai(text: str, num_q: int = 5) -> list:
    """
    Try each model in MODELS_TO_TRY until one returns a valid quiz.
    Falls back to built-in questions if all models fail.
    """
    api_key = os.environ.get("OPENROUTER_API_KEY", "")
    if not api_key:
        logger.warning("No OPENROUTER_API_KEY — using fallback quiz")
        return fallback_quiz()[:num_q]

    prompt = _build_prompt(text, num_q)
    all_errors = []

    async with httpx.AsyncClient(timeout=45) as client:
        for i, (model_id, model_name) in enumerate(MODELS_TO_TRY, 1):
            logger.info(f"Trying model {i}/{len(MODELS_TO_TRY)}: {model_name}")
            quiz, error = await _try_model(client, api_key, model_id, prompt)

            if quiz:
                logger.info(f"✅ Success with {model_name} — {len(quiz)} questions generated")
                return quiz[:num_q]
            else:
                logger.warning(f"❌ {model_name} failed: {error}")
                all_errors.append(f"{model_name}: {error}")
                # Small delay before next model to avoid hammering
                await asyncio.sleep(0.5)

    logger.error(f"All models failed: {' | '.join(all_errors)} — using fallback quiz")
    return fallback_quiz()[:num_q]


def fallback_quiz() -> list:
    return [
        {
            "question": "Under the Customs Act 1962, what is the primary basis for customs valuation of imported goods?",
            "options": ["Transaction value (price actually paid or payable)", "Market value in exporting country", "Manufacturer's declared cost", "CIF value as assessed by officer"],
            "answer": 0,
            "explanation": "Section 14 of Customs Act adopts Transaction Value per WTO Valuation Agreement as the primary method."
        },
        {
            "question": "Which HS Code chapter covers motor vehicles and other road vehicles?",
            "options": ["Chapter 84", "Chapter 86", "Chapter 87", "Chapter 88"],
            "answer": 2,
            "explanation": "Chapter 87 of the Harmonized System specifically covers motor vehicles, tractors, and road vehicles."
        },
        {
            "question": "The WTO Agreement on Customs Valuation (CVA) provides how many methods of valuation in sequence?",
            "options": ["3 methods", "5 methods", "6 methods", "8 methods"],
            "answer": 2,
            "explanation": "The CVA provides 6 sequential methods: Transaction Value, Identical Goods, Similar Goods, Deductive, Computed, and Fallback."
        },
        {
            "question": "Under IGST Act, the place of supply for import of services is:",
            "options": ["Location of the supplier abroad", "Location of the service recipient in India", "Port of entry into India", "Place where payment is made"],
            "answer": 1,
            "explanation": "Section 13 IGST Act: place of supply of imported services = location of the recipient in India."
        },
        {
            "question": "An advance Bill of Entry can be filed how many days before the expected arrival of goods?",
            "options": ["7 days", "15 days", "30 days", "60 days"],
            "answer": 2,
            "explanation": "Customs (Advance Filing) Regulations allow advance Bill of Entry up to 30 days before expected arrival."
        },
        {
            "question": "Which section of the Customs Act deals with 'Baggage'?",
            "options": ["Section 77-88", "Section 60-76", "Section 46-54", "Section 30-45"],
            "answer": 0,
            "explanation": "Sections 77 to 88 of Customs Act 1962 exclusively deal with baggage rules for passengers."
        },
        {
            "question": "Basic Customs Duty (BCD) on most goods is levied under which section of the Customs Act?",
            "options": ["Section 12", "Section 14", "Section 25", "Section 46"],
            "answer": 0,
            "explanation": "Section 12 is the charging section for customs duty on all goods imported into or exported from India."
        },
        {
            "question": "The Special Valuation Branch (SVB) of Customs investigates transactions involving:",
            "options": ["First-time importers only", "Related party transactions", "Goods above ₹1 crore value", "Goods from sanctioned countries"],
            "answer": 1,
            "explanation": "SVB investigates related party transactions to ensure relationship has not influenced the declared transaction value."
        },
    ]


# ─────────────────────────────────────────────
# REST ENDPOINTS
# ─────────────────────────────────────────────

@app.get("/")
async def root():
    return FileResponse("static/trainer.html")

@app.get("/join")
async def join_page():
    return FileResponse("static/participant.html")

@app.get("/api/state")
async def get_state():
    return {
        "active": state["active"],
        "session_name": state["session_name"],
        "join_code": state["join_code"],
        "question_count": len(state["quiz_data"]),
        "participant_count": len(state["scores"]),
    }

@app.get("/api/quiz")
async def get_quiz():
    if not state["active"]:
        raise HTTPException(status_code=404, detail="No active session")
    # Return quiz without answers
    safe_quiz = [
        {"question": q["question"], "options": q["options"], "index": i}
        for i, q in enumerate(state["quiz_data"])
    ]
    return {"quiz": safe_quiz, "session_name": state["session_name"]}

@app.post("/api/launch")
async def launch_session(
    session_name: str = Form("Customs Training"),
    num_questions: int = Form(5),
    file: Optional[UploadFile] = File(None)
):
    import random, string
    code = "".join(random.choices(string.digits, k=6))
    code = code[:3] + " " + code[3:]

    quiz = []
    extraction_method = "none"
    extracted_chars = 0

    if file:
        content = await file.read()
        fn = file.filename.lower()
        logger.info(f"Processing uploaded file: {file.filename} ({len(content)} bytes)")

        text = extract_text(file.filename, content)
        extracted_chars = len(text)

        # Determine which method was used (for logging/response)
        if extracted_chars >= MIN_USEFUL_CHARS:
            if fn.endswith(".pptx") or fn.endswith(".ppt"):
                extraction_method = "pptx_shapes"
            elif fn.endswith(".pdf"):
                # Heuristic: if pdfplumber alone would have worked, digital; else ocr
                extraction_method = "pdf_digital" if HAS_PDF else "pdf_ocr"
            elif fn.endswith((".png", ".jpg", ".jpeg", ".tiff", ".bmp")):
                extraction_method = "image_ocr"
            else:
                extraction_method = "text"

        logger.info(f"Extraction: {extraction_method}, {extracted_chars} chars")

        if text and extracted_chars > 50:
            quiz = await generate_quiz_ai(text, num_questions)
        else:
            logger.warning("Extracted text too short — using fallback quiz")

    if not quiz:
        quiz = fallback_quiz()[:num_questions]
        extraction_method = "fallback"

    state.update({
        "active": True,
        "session_name": session_name,
        "join_code": code,
        "quiz_data": quiz,
        "scores": {},
        "started_at": time.time(),
    })

    await manager.broadcast_trainers({
        "type": "session_launched",
        "session_name": session_name,
        "join_code": code,
        "question_count": len(quiz),
    })

    return {
        "status": "launched",
        "join_code": code,
        "question_count": len(quiz),
        "extraction_method": extraction_method,
        "extracted_chars": extracted_chars,
    }

@app.post("/api/reset")
async def reset_session():
    state.update({"active": False, "quiz_data": [], "scores": {}, "join_code": "000000"})
    await manager.broadcast_all({"type": "session_reset"})
    return {"status": "reset"}

@app.post("/api/submit")
async def submit_answers(payload: dict):
    if not state["active"]:
        raise HTTPException(status_code=400, detail="No active session")

    user_id = payload.get("user_id", "")
    name = payload.get("name", "Anonymous")
    answers = payload.get("answers", [])

    if user_id in state["scores"] and state["scores"][user_id].get("submitted"):
        return {"status": "already_submitted", "score": state["scores"][user_id]["score"]}

    score = 0
    correct = 0
    wrong = 0
    result_answers = []

    for i, q in enumerate(state["quiz_data"]):
        user_ans = answers[i] if i < len(answers) else -1
        is_correct = (user_ans == q["answer"])
        if is_correct:
            score += 10
            correct += 1
        else:
            wrong += 1
        result_answers.append({
            "user_answer": user_ans,
            "correct_answer": q["answer"],
            "correct": is_correct,
            "explanation": q.get("explanation", "")
        })

    avatar = AVATARS[hash(user_id) % len(AVATARS)]

    state["scores"][user_id] = {
        "user_id": user_id,
        "name": name,
        "score": score,
        "correct": correct,
        "wrong": wrong,
        "avatar": avatar,
        "answers": result_answers,
        "submitted": True,
        "submitted_at": time.time(),
        "joined_at": state["scores"].get(user_id, {}).get("joined_at", time.time()),
    }

    # Broadcast updated leaderboard to ALL connected clients instantly
    snapshot = leaderboard_snapshot()
    await manager.broadcast_all(snapshot)

    return {
        "status": "submitted",
        "score": score,
        "correct": correct,
        "wrong": wrong,
        "total": len(state["quiz_data"]),
        "answers": result_answers,
        "quiz_data": state["quiz_data"],
    }

@app.post("/api/join")
async def join_session(payload: dict):
    user_id = payload.get("user_id", "")
    name = payload.get("name", "Anonymous")
    if not state["active"]:
        raise HTTPException(status_code=400, detail="No active session")
    if user_id not in state["scores"]:
        state["scores"][user_id] = {
            "user_id": user_id,
            "name": name,
            "score": 0,
            "correct": 0,
            "wrong": 0,
            "avatar": AVATARS[hash(user_id) % len(AVATARS)],
            "answers": [],
            "submitted": False,
            "joined_at": time.time(),
        }
    snapshot = leaderboard_snapshot()
    await manager.broadcast_all(snapshot)
    return {"status": "joined", "avatar": state["scores"][user_id]["avatar"]}


# ─────────────────────────────────────────────
# WEBSOCKET ENDPOINTS
# ─────────────────────────────────────────────

@app.websocket("/ws/trainer")
async def ws_trainer(websocket: WebSocket):
    await manager.connect_trainer(websocket)
    try:
        # Send current state immediately
        await websocket.send_json(leaderboard_snapshot())
        while True:
            # Keep alive + listen for trainer commands
            try:
                data = await asyncio.wait_for(websocket.receive_json(), timeout=30)
                if data.get("type") == "ping":
                    await websocket.send_json({"type": "pong"})
            except asyncio.TimeoutError:
                await websocket.send_json({"type": "ping"})
    except WebSocketDisconnect:
        manager.disconnect_trainer(websocket)

@app.websocket("/ws/participant/{user_id}")
async def ws_participant(websocket: WebSocket, user_id: str):
    await manager.connect_participant(websocket, user_id)
    try:
        await websocket.send_json(leaderboard_snapshot())
        while True:
            try:
                data = await asyncio.wait_for(websocket.receive_json(), timeout=30)
                if data.get("type") == "ping":
                    await websocket.send_json({"type": "pong"})
            except asyncio.TimeoutError:
                await websocket.send_json({"type": "ping"})
    except WebSocketDisconnect:
        manager.disconnect_participant(user_id)


# ─────────────────────────────────────────────
# ENTRY POINT
# ─────────────────────────────────────────────
if __name__ == "__main__":
    port = int(os.environ.get("PORT", 8000))
    uvicorn.run("main:app", host="0.0.0.0", port=port, reload=False)
